import requests
import json
import re
import plugins
from bridge.reply import Reply, ReplyType
from bridge.context import ContextType
from plugins import *
from common.log import logger
import os
from docx import Document
import markdown
import tiktoken
import jieba
import fitz
from openpyxl import load_workbook
import csv
from bs4 import BeautifulSoup
from pptx import Presentation
import base64

EXTENSION_TO_TYPE = {
    'pdf': 'pdf',
    'doc': 'docx', 'docx': 'docx',
    'md': 'md',
    'txt': 'txt',
    'xls': 'excel', 'xlsx': 'excel',
    'csv': 'csv',
    'html': 'html', 'htm': 'html',
    'ppt': 'ppt', 'pptx': 'ppt'
}

@plugins.register(
    name="sum4all",
    desire_priority=2,
    desc="A plugin for summarizing all things",
    version="0.3.6",
    author="fatwang2",
)


class sum4all(Plugin):
    def __init__(self):
        super().__init__()
        try:
            curdir = os.path.dirname(__file__)
            config_path = os.path.join(curdir, "config.json")
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    self.config = json.load(f)
            else:
                # 使用父类的方法来加载配置
                self.config = super().load_config()
                
                if not self.config:
                    raise Exception("config.json not found")
            # 设置事件处理函数
            self.handlers[Event.ON_HANDLE_CONTEXT] = self.on_handle_context
            # 从配置中提取所需的设置
            self.sum_service = self.config.get("sum_service","")
            self.bibigpt_key = self.config.get("bibigpt_key","")
            self.outputLanguage = self.config.get("outputLanguage","zh-CN")
            self.group_sharing = self.config.get("group_sharing","true")
            self.opensum_key = self.config.get("opensum_key","")
            self.open_ai_api_key = self.config.get("open_ai_api_key","")
            self.model = self.config.get("model","gpt-3.5-turbo")
            self.open_ai_api_base = self.config.get("open_ai_api_base","https://api.openai.com/v1")
            self.prompt = self.config.get("prompt","你是一个新闻专家，我会给你发一些网页内容，请你用简单明了的语言做总结。格式如下：📌总结\n一句话讲清楚整篇文章的核心观点，控制在30字左右。\n\n💡要点\n用数字序号列出来3-5个文章的核心内容，尽量使用emoji让你的表达更生动，请注意输出的内容不要有两个转义符")
            self.search_prompt = self.config.get("search_prompt","你是一个信息检索专家，请你用简单明了的语言，对你收到的内容做总结。尽量使用emoji让你的表达更生动")
            self.sum4all_key = self.config.get("sum4all_key","")
            self.search_sum = self.config.get("search_sum","")
            self.file_sum = self.config.get("file_sum","")
            self.image_sum = self.config.get("image_sum","")
            self.perplexity_key = self.config.get("perplexity_key","")
            self.search_service = self.config.get("search_service","")            
                
            # 初始化成功日志
            logger.info("[sum4all] inited.")
        except Exception as e:
            # 初始化失败日志
            logger.warn(f"sum4all init failed: {e}")
    def on_handle_context(self, e_context: EventContext):
        context = e_context["context"]
        if context.type not in [ContextType.TEXT, ContextType.SHARING,ContextType.FILE,ContextType.IMAGE]:
            return
        content = context.content
        isgroup = e_context["context"].get("isgroup", False)

        url_match = re.match('https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+', content)
        unsupported_urls = re.search(r'.*finder\.video\.qq\.com.*|.*support\.weixin\.qq\.com/update.*|.*support\.weixin\.qq\.com/security.*|.*mp\.weixin\.qq\.com/mp/waerrpage.*', content)

            # 检查输入是否以"搜" 开头
        if content.startswith("搜") and self.search_sum:
            # Call new function to handle search operation
            self.call_service(content, e_context, "search")
            return
        if context.type == ContextType.FILE:
            logger.info("on_handle_context: 处理上下文开始")
            context.get("msg").prepare()
            file_path = context.content
            logger.info(f"on_handle_context: 获取到文件路径 {file_path}")
            # 检查是否应该进行文件总结
            if self.file_sum:
                content = self.extract_content(file_path)  # 提取内容
                self.handle_openai_file(content, e_context)
            else:
                logger.info("文件总结功能已禁用，不对文件内容进行处理")
        elif context.type == ContextType.IMAGE:
            logger.info("on_handle_context: 开始处理图片")
            context.get("msg").prepare()
            image_path = context.content
            logger.info(f"on_handle_context: 获取到图片路径 {image_path}")
            # 检查是否应该进行图片总结
            if self.image_sum:
                self.handle_openai_image(image_path, e_context)
            else:
                logger.info("图片总结功能已禁用，不对图片内容进行处理")
        elif context.type == ContextType.SHARING:  #匹配卡片分享
            if unsupported_urls:  #匹配不支持总结的卡片
                if isgroup:  ##群聊中忽略
                    return
                else:  ##私聊回复不支持
                    logger.info("[sum4all] Unsupported URL : %s", content)
                    reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                    e_context["reply"] = reply
                    e_context.action = EventAction.BREAK_PASS
            else:  #匹配支持总结的卡片
                if isgroup:  #处理群聊总结
                    if self.group_sharing:  #group_sharing = True进行总结，False则忽略。
                        logger.info("[sum4all] Summary URL : %s", content)
                        self.call_service(content, e_context, "sum")
                        return
                    else:
                        return
                else:  #处理私聊总结
                    logger.info("[sum4all] Summary URL : %s", content)
                    self.call_service(content, e_context, "sum")
                    return
        elif url_match: #匹配URL链接
            if unsupported_urls:  #匹配不支持总结的网址
                logger.info("[sum4all] Unsupported URL : %s", content)
                reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                e_context["reply"] = reply
                e_context.action = EventAction.BREAK_PASS
            else:
                logger.info("[sum4all] Summary URL : %s", content)
                self.call_service(content, e_context, "sum")
                return
    def call_service(self, content, e_context, service_type):
        if service_type == "search":
            if self.search_service == "sum4all":
                self.handle_search(content, e_context)
            elif self.search_service == "perplexity":
                self.handle_perplexity(content, e_context)
        elif service_type == "sum":
            if self.sum_service == "bibigpt":
                self.handle_bibigpt(content, e_context)
            elif self.sum_service == "openai":
                self.handle_openai(content, e_context)
            elif self.sum_service == "opensum":
                self.handle_opensum(content, e_context)
            elif self.sum_service == "sum4all":
                self.handle_sum4all(content, e_context)
    
    def short_url(self, long_url):
        url = "https://short.fatwang2.com"
        payload = {
            "url": long_url
        }        
        headers = {'Content-Type': "application/json"}
        response = requests.request("POST", url, json=payload, headers=headers)
        if response.status_code == 200:
            res_data = response.json()
            # 直接从返回的 JSON 中获取短链接
            short_url = res_data.get('shorturl', None)  
            
            if short_url:
                return short_url
        return None
    def handle_openai(self, content, e_context):
        meta = None      
        headers = {
            'Content-Type': 'application/json',
            'WebPilot-Friend-UID': 'fatwang2'
        }
        payload = json.dumps({"link": content})
        try:
            api_url = "https://gpts.webpilot.ai/api/visit-web"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            meta= data.get('content','content not available')  # 获取data字段                

        except requests.exceptions.RequestException as e:
            meta = f"An error occurred: {e}"          

        # 如果meta获取成功，发送请求到OpenAI
        if meta:
            try:
                headers = {
                    'Content-Type': 'application/json',
                    'Authorization': f'Bearer {self.open_ai_api_key}'  # 使用你的OpenAI API密钥
                }
                data = {
                    "model": self.model, 
                    "messages": [
                        {"role": "system", "content": self.prompt},
                        {"role": "user", "content": meta}
                    ]
                }
            
                response = requests.post(f"{self.open_ai_api_base}/chat/completions", headers=headers, data=json.dumps(data))
                response.raise_for_status()

                # 处理响应数据
                response_data = response.json()
                # 这里可以根据你的需要处理响应数据
                # 解析 JSON 并获取 content
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    first_choice = response_data["choices"][0]
                    if "message" in first_choice and "content" in first_choice["message"]:
                        content = first_choice["message"]["content"]
                    else:
                        print("Content not found in the response")
                else:
                    print("No choices available in the response")
            except requests.exceptions.RequestException as e:
                # 处理可能出现的错误
                logger.error(f"Error calling OpenAI API: {e}")
            reply = Reply()
            reply.type = ReplyType.TEXT
            reply.content = f"{content}"            
            e_context["reply"] = reply
            e_context.action = EventAction.BREAK_PASS
    def handle_sum4all(self, content, e_context):
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.sum4all_key}'
        }
        payload = json.dumps({
            "link": content,
            "prompt": self.prompt
        })
        additional_content = ""  # 在 try 块之前初始化 additional_content

        try:
            api_url = "https://ai.sum4all.site"
            response = requests.post(api_url, headers=headers, data=payload)
            response.raise_for_status()
            response_data = response.json()  # 解析响应的 JSON 数据
            if response_data.get("success"):
                content = response_data["content"].replace("\\n", "\n")  # 替换 \\n 为 \n

                # 新增加的部分，用于解析 meta 数据
                meta = response_data.get("meta", {})  # 如果没有 meta 数据，则默认为空字典
                title = meta.get("og:title", "")  # 获取 og:title，如果没有则默认为空字符串
                # 只有当 title 非空时，才加入到回复中
                if title:
                    additional_content += f"{title}\n\n"
                reply_content = additional_content + content  # 将内容加入回复
                
            else:
                reply_content = "Content not found or error in response"

        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling new combined api: {e}")
            reply_content = f"An error occurred: {e}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = reply_content            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_bibigpt(self, content, e_context):    
        headers = {
            'Content-Type': 'application/json'
        }
        payload_params = {
            "url": content,
            "includeDetail": False,
            "promptConfig": {
                "outputLanguage": self.outputLanguage
            }
        }

        payload = json.dumps(payload_params)           
        try:
            api_url = f"https://bibigpt.co/api/open/{self.bibigpt_key}"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_original = data.get('summary', 'Summary not available')
            html_url = data.get('htmlUrl', 'HTML URL not available')
            # 获取短链接
            short_url = self.short_url(html_url) 
            
            # 如果获取短链接失败，使用 html_url
            if short_url is None:
                short_url = html_url if html_url != 'HTML URL not available' else 'URL not available'
            
            # 移除 "##摘要"、"## 亮点" 和 "-"
            summary = summary_original.split("详细版（支持对话追问）")[0].replace("## 摘要\n", "📌总结：").replace("## 亮点\n", "").replace("- ", "")
        except requests.exceptions.RequestException as e:
            reply = f"An error occurred: {e}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS


    def handle_opensum(self, content, e_context):
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.opensum_key}'
        }
        payload = json.dumps({"link": content})
        try:
            api_url = "https://read.thinkwx.com/api/v1/article/summary"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_data = data.get('data', {})  # 获取data字段                
            summary_original = summary_data.get('summary', 'Summary not available')
            # 使用正则表达式提取URL
            url_pattern = r'https:\/\/[^\s]+'
            match = re.search(url_pattern, summary_original)
            html_url = match.group(0) if match else 'HTML URL not available'            
            # 获取短链接
            short_url = self.short_url(html_url) if match else html_url
            # 用于移除摘要中的URL及其后的所有内容
            url_pattern_remove = r'https:\/\/[^\s]+[\s\S]*'
            summary = re.sub(url_pattern_remove, '', summary_original).strip()        

        except requests.exceptions.RequestException as e:
            summary = f"An error occurred: {e}"
            short_url = 'URL not available'
        
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS    
    def handle_search(self, content, e_context):
        
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.sum4all_key}'
        }
        payload = json.dumps({
            "ur": content,
            "prompt": self.search_prompt
        })
        try:
            api_url = "https://ai.sum4all.site"
            response = requests.post(api_url, headers=headers, data=payload)
            response.raise_for_status()
            response_data = response.json()  # 解析响应的 JSON 数据
            if response_data.get("success"):
                content = response_data["content"].replace("\\n", "\n")  # 替换 \\n 为 \n
                reply_content = content  # 将内容加入回复

                # 解析 meta 数据
                meta = response_data.get("meta", {})  # 如果没有 meta 数据，则默认为空字典
                title = meta.get("og:title", "")  # 获取 og:title，如果没有则默认为空字符串
                og_url = meta.get("og:url", "")  # 获取 og:url，如果没有则默认为空字符串
                # 打印 title 和 og_url 以调试
                print("Title:", title)
                print("Original URL:", og_url)                
                # 只有当 title 和 url 非空时，才加入到回复中
                if title:
                    reply_content += f"\n\n参考文章：{title}"
                if og_url:
                    short_url = self.short_url(og_url)  # 获取短链接
                    reply_content += f"\n\n参考链接：{short_url}"                

            else:
                content = "Content not found or error in response"

        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling new combined api: {e}")
            reply_content = f"An error occurred: {e}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = reply_content            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_perplexity(self, content, e_context):
        
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.perplexity_key}'
        }
        data = {
            "model": "pplx-7b-online", 
            "messages": [
                {"role": "system", "content": self.search_prompt},
                {"role": "user", "content": content}
        ]
        }
        try:
            api_url = "https://api.perplexity.ai/chat/completions"
            response = requests.post(api_url, headers=headers, json=data)
            response.raise_for_status()
            # 处理响应数据
            response_data = response.json()
            # 这里可以根据你的需要处理响应数据
            # 解析 JSON 并获取 content
            if "choices" in response_data and len(response_data["choices"]) > 0:
                first_choice = response_data["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    content = first_choice["message"]["content"]
                else:
                    print("Content not found in the response")
            else:
                print("No choices available in the response")
        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling perplexity: {e}")
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{content}"            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def get_help_text(self, **kwargs):
        help_text = "输入url/分享链接/搜索关键词，直接为你总结\n"
        return help_text
    def handle_openai_file(self, content, e_context):
        logger.info("handle_openai_file: 向OpenAI发送内容总结请求")
        # 根据sum_service的值选择API密钥和基础URL
        if self.sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
        elif self.sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
        else:
            logger.error(f"未知的sum_service配置: {self.sum_service}")
            return

        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        data = {
            "model": model, 
            "messages": [
                {"role": "system", "content": self.prompt},
                {"role": "user", "content": content}
            ]
        }
        try:
            response = requests.post(f"{api_base}/chat/completions", headers=headers, data=json.dumps(data))
            response.raise_for_status()
            response_data = response.json()

            # 解析 JSON 并获取 content
            if "choices" in response_data and len(response_data["choices"]) > 0:
                first_choice = response_data["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    response_content = first_choice["message"]["content"].strip()  # 获取响应内容
                    logger.info(f"OpenAI API response content")  # 记录响应内容
                    reply_content = response_content.replace("\\n", "\n")  # 替换 \\n 为 \n
                else:
                    logger.error("Content not found in the response")
                    reply_content = "Content not found in the OpenAI API response"
            else:
                logger.error("No choices available in the response")
                reply_content = "No choices available in the OpenAI API response"

        except requests.exceptions.RequestException as e:
            logger.error(f"Error calling OpenAI API: {e}")
            reply_content = f"An error occurred while calling OpenAI API: {e}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = reply_content  # 设置响应内容到回复对象
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS


    def read_pdf(self, file_path):
        logger.info(f"开始读取PDF文件：{file_path}")
        doc = fitz.open(file_path)
        content = ' '.join([page.get_text() for page in doc])
        logger.info(f"PDF文件读取完成：{file_path}")

        return content

    def read_word(self, file_path):
        doc = Document(file_path)
        return ' '.join([p.text for p in doc.paragraphs])

    def read_markdown(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
            return markdown.markdown(md_content)

    def read_excel(self, file_path):
        workbook = load_workbook(file_path)
        content = ''
        for sheet in workbook:
            for row in sheet.iter_rows():
                content += ' '.join([str(cell.value) for cell in row])
                content += '\n'
        return content

    def read_txt(self, file_path):
        logger.debug(f"开始读取TXT文件: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            logger.debug(f"TXT文件读取完成: {file_path}")
            logger.debug("TXT文件内容的前50个字符：")
            logger.debug(content[:50])  # 打印文件内容的前50个字符
            return content
        except Exception as e:
            logger.error(f"读取TXT文件时出错: {file_path}，错误信息: {str(e)}")
            return ""

    def read_csv(self, file_path):
        content = ''
        with open(file_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                content += ' '.join(row) + '\n'
        return content

    def num_tokens_from_string(self, text):
        try:
            encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
        except KeyError:
            logger.debug(f"Warning: model not found. Using cl100k_base encoding.")
            encoding = tiktoken.get_encoding("cl100k_base")
        return len(encoding.encode(text))

    def read_html(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            return soup.get_text()

    def read_ppt(self, file_path):
        presentation = Presentation(file_path)
        content = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + '\n'
        return content

    def split_text_chinese(self, text, overlap_tokens=500):
        tokens = jieba.cut(text)
        segments = []
        segment_text = ""
        for token in tokens:
            temp_segment_text = segment_text + token
            temp_segment_tokens_count = self.num_tokens_from_string(temp_segment_text)
            if temp_segment_tokens_count >= self.max_tokens:
                segments.append(segment_text)
                previous_segment_text = segment_text
                segment_text = previous_segment_text[-overlap_tokens:] + token if overlap_tokens > 0 else token
            else:
                segment_text = temp_segment_text

        if segment_text:
            segments.append(segment_text)
        logger.debug(f"分段文本: {segments}")
        return segments
    def extract_content(self, file_path):
        logger.info(f"extract_content: 提取文件内容，文件路径: {file_path}")
    
        file_extension = os.path.splitext(file_path)[1][1:].lower()
        logger.info(f"extract_content: 文件类型为 {file_extension}")
    
        file_type = EXTENSION_TO_TYPE.get(file_extension)

        if not file_type:
            logger.error(f"不支持的文件扩展名: {file_extension}")
            return None

        read_func = {
            'pdf': self.read_pdf,
            'docx': self.read_word,
            'md': self.read_markdown,
            'txt': self.read_txt,
            'excel': self.read_excel,
            'csv': self.read_csv,
            'html': self.read_html,
            'ppt': self.read_ppt
        }.get(file_type)

        if not read_func:
            logger.error(f"不支持的文件类型: {file_type}")
            return None
        logger.info("extract_content: 文件内容提取完成")
        return read_func(file_path)
    # Function to encode the image
    def encode_image(image_path):
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    # Function to handle OpenAI image processing
    def handle_openai_image(self, image_path, e_context):
        logger.info("handle_openai_image_response: 解析OpenAI图像处理API的响应")

        # Getting the base64 string
        base64_image = encode_image(image_path)

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.open_ai_api_key}"
        }

        payload = {
            "model": "gpt-4-vision-preview",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "图片讲了什么?"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 300
        }

        try:
            response = requests.post(f"{self.open_ai_api_base}/chat/completions", headers=headers, json=payload)
            response.raise_for_status()  # 增加对HTTP错误的检查
            response_json = response.json()  # 定义response_json
            # 确保响应中有 'choices' 键并且至少有一个元素
            if "choices" in response_json and len(response_json["choices"]) > 0:
                first_choice = response_json["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    # 从响应中提取 'content'
                    response_content = first_choice["message"]["content"].strip()
                    logger.info("OpenAI API response content")  # 记录响应内容
                    reply_content = response_content
                else:
                    logger.error("Content not found in the response")
                    reply_content = "Content not found in the OpenAI API response"
            else:
                logger.error("No choices available in the response")
                reply_content = "No choices available in the OpenAI API response"
        except Exception as e:
            logger.error(f"Error processing OpenAI API response: {e}")
            reply_content = f"An error occurred while processing OpenAI API response: {e}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = reply_content  # 设置响应内容到回复对象
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS